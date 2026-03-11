import os
import json
import requests
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import google.generativeai as genai
from dotenv import load_dotenv
from io import BytesIO
import datetime

# Load environment variables
load_dotenv()

# Configure API keys
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
UNSPLASH_ACCESS_KEY = os.getenv("UNSPLASH_ACCESS_KEY")

if GEMINI_API_KEY:
    genai.configure(api_key=GEMINI_API_KEY)

# Set Streamlit page configuration
st.set_page_config(
    page_title="AI Presentation Builder",
    page_icon="📊",
    layout="wide"
)

# Custom CSS for better Streamlit UI
st.markdown("""
    <style>
    .stButton > button {
        width: 100%;
        background: linear-gradient(90deg, #4776E6 0%, #8E54E9 100%);
        color: white;
        font-weight: bold;
        border: none;
        padding: 10px 20px;
        border-radius: 5px;
    }
    .stDownloadButton > button {
        background: linear-gradient(90deg, #11998e 0%, #38ef7d 100%);
        color: white;
        font-weight: bold;
    }
    .title-text {
        text-align: center;
        color: #4776E6;
    }
    </style>
""", unsafe_allow_html=True)

st.markdown("<h1 class='title-text'>📊 AI Powered Presentation Builder</h1>", unsafe_allow_html=True)
st.markdown("Generate beautiful PowerPoint presentations using AI")

# Function to get available models
@st.cache_data
def get_available_models():
    """Get list of available Gemini models."""
    try:
        models = genai.list_models()
        available_models = []
        for model in models:
            if 'generateContent' in model.supported_generation_methods:
                # Clean up model name
                model_name = model.name.replace('models/', '')
                available_models.append(model_name)
        return available_models
    except Exception as e:
        st.error(f"Error fetching models: {e}")
        return []

# Check available models if API key is configured
available_models = []
if GEMINI_API_KEY:
    with st.spinner("Checking available models..."):
        available_models = get_available_models()

# Sidebar for configuration
with st.sidebar:
    st.header("⚙️ Configuration")
    
    # Show API status
    if GEMINI_API_KEY:
        st.success("✅ Gemini API Connected")
        if available_models:
            st.info(f"📚 {len(available_models)} models available")
    else:
        st.error("❌ Gemini API Key missing")
    
    num_slides = st.slider("Number of slides", 3, 15, 5)
    
    # Model selection (if models are available)
    if available_models:
        default_model = next((m for m in available_models if 'flash' in m), available_models[0])
        selected_model = st.selectbox(
            "AI Model",
            available_models,
            index=available_models.index(default_model) if default_model in available_models else 0
        )
    else:
        selected_model = "gemini-pro"  # Fallback
        st.warning("Using default model")
    
    presentation_style = st.selectbox(
        "Presentation style",
        ["Professional", "Educational", "Creative", "Minimal", "Corporate"]
    )
    
    # Color scheme based on style
    color_schemes = {
        "Professional": {
            "primary": RGBColor(41, 128, 185),
            "secondary": RGBColor(52, 73, 94),
            "accent": RGBColor(46, 204, 113)
        },
        "Educational": {
            "primary": RGBColor(39, 174, 96),
            "secondary": RGBColor(241, 196, 15),
            "accent": RGBColor(155, 89, 182)
        },
        "Creative": {
            "primary": RGBColor(155, 89, 182),
            "secondary": RGBColor(230, 126, 34),
            "accent": RGBColor(52, 152, 219)
        },
        "Minimal": {
            "primary": RGBColor(52, 73, 94),
            "secondary": RGBColor(149, 165, 166),
            "accent": RGBColor(127, 140, 141)
        },
        "Corporate": {
            "primary": RGBColor(44, 62, 80),
            "secondary": RGBColor(52, 152, 219),
            "accent": RGBColor(46, 204, 113)
        }
    }
    
    colors = color_schemes[presentation_style]
    
    st.divider()
    include_images = st.checkbox("Include images", value=True)
    include_page_numbers = st.checkbox("Include page numbers", value=True)

# Main content area
col1, col2 = st.columns([2, 1])

with col1:
    topic = st.text_input(
        "**Enter presentation topic**",
        placeholder="e.g., Machine Learning, Climate Change, Web Development"
    )

with col2:
    generate_btn = st.button("🚀 Generate Presentation", use_container_width=True)

# Progress tracking
progress_placeholder = st.empty()
status_placeholder = st.empty()

def fetch_image_from_unsplash(query: str) -> str:
    """Fetch a relevant image from Unsplash API."""
    if not UNSPLASH_ACCESS_KEY or not include_images:
        return None
    
    try:
        url = f"https://api.unsplash.com/search/photos"
        params = {
            "query": query,
            "per_page": 1,
            "orientation": "landscape",
            "content_filter": "high"
        }
        headers = {"Authorization": f"Client-ID {UNSPLASH_ACCESS_KEY}"}
        
        response = requests.get(url, params=params, headers=headers, timeout=10)
        if response.status_code == 200:
            data = response.json()
            if data["results"]:
                # Get high quality image
                return data["results"][0]["urls"]["regular"] + "&w=1200&h=800&fit=crop"
    except Exception as e:
        st.warning(f"Could not fetch image: {e}")
    
    return None

def generate_presentation_content(topic: str, num_slides: int, style: str, model_name: str) -> dict:
    """Generate presentation content using Gemini API."""
    
    prompt = f"""Generate a comprehensive presentation outline for the topic: '{topic}'
    
    Style: {style}
    Number of slides: {num_slides}
    
    Return ONLY a valid JSON object with this exact structure. No other text or explanation:
    {{
        "title": "Presentation Title",
        "subtitle": "Engaging subtitle for the presentation",
        "slides": [
            {{
                "slide_number": 1,
                "title": "Slide Title",
                "content": "Main point 1|Main point 2|Main point 3|Main point 4",
                "image_query": "relevant search term for images"
            }}
        ]
    }}
    
    Requirements:
    - Create exactly {num_slides} slides
    - First slide should introduce the topic
    - Last slide should be a conclusion
    - Each slide should have 3-4 bullet points separated by |
    - Make image queries specific and visual
    """
    
    try:
        # Use the selected model
        model = genai.GenerativeModel(model_name)
        response = model.generate_content(
            prompt,
            generation_config=genai.types.GenerationConfig(
                temperature=0.7,
                max_output_tokens=2048,
                top_p=0.8,
                top_k=40
            )
        )
        
        # Extract JSON from response
        response_text = response.text
        
        # Clean up response (remove markdown if present)
        if "```json" in response_text:
            response_text = response_text.split("```json")[1].split("```")[0]
        elif "```" in response_text:
            response_text = response_text.split("```")[1].split("```")[0]
        
        # Try to find JSON in the response
        start_idx = response_text.find('{')
        end_idx = response_text.rfind('}') + 1
        
        if start_idx != -1 and end_idx > start_idx:
            json_str = response_text[start_idx:end_idx]
            return json.loads(json_str)
        else:
            raise ValueError("No JSON found in response")
    
    except Exception as e:
        st.error(f"Error generating content: {e}")
        return None

def create_beautiful_title_slide(prs, title, subtitle, colors):
    """Create a beautiful title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors["primary"]
    
    # Add decorative elements
    left = Inches(0)
    top = Inches(0)
    width = Inches(10)
    height = Inches(7.5)
    
    # Add diagonal accent
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = colors["secondary"]
    accent.fill.transparency = 0.8
    accent.rotation = 15
    accent.line.fill.background()
    
    # Add title
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(2)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    if subtitle:
        left = Inches(0.5)
        top = Inches(4.5)
        width = Inches(9)
        height = Inches(1)
        
        subtitle_box = slide.shapes.add_textbox(left, top, width, height)
        subtitle_frame = subtitle_box.text_frame
        
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        p.font.italic = True
    
    # Add date
    left = Inches(0.5)
    top = Inches(6.5)
    width = Inches(9)
    height = Inches(0.5)
    
    date_box = slide.shapes.add_textbox(left, top, width, height)
    date_frame = date_box.text_frame
    p = date_frame.paragraphs[0]
    p.text = datetime.datetime.now().strftime("%B %d, %Y")
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    p.font.italic = True

def create_beautiful_content_slide(prs, slide_data, colors, slide_number, total_slides):
    """Create a beautiful content slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # White background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Add header bar
    left = Inches(0)
    top = Inches(0)
    width = Inches(10)
    height = Inches(1.2)
    
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    header.fill.solid()
    header.fill.fore_color.rgb = colors["primary"]
    header.line.fill.background()
    
    # Add accent line
    left = Inches(0)
    top = Inches(1.2)
    width = Inches(10)
    height = Inches(0.1)
    
    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = colors["secondary"]
    accent_line.line.fill.background()
    
    # Add slide title
    left = Inches(0.3)
    top = Inches(0.2)
    width = Inches(9.4)
    height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = slide_data.get("title", f"Slide {slide_number}")
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Add content area
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(5.5)
    height = Inches(4.5)
    
    content_box = slide.shapes.add_textbox(left, top, width, height)
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    # Add bullet points
    content = slide_data.get("content", "").split("|")
    for i, point in enumerate(content):
        if point.strip():
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            p.text = point.strip()
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(44, 62, 80)
            p.level = 0
            p.space_before = Pt(12)
            p.space_after = Pt(6)
            p.bullet = True
    
    # Add image if available
    if slide_data.get("image_query") and UNSPLASH_ACCESS_KEY and include_images:
        image_url = fetch_image_from_unsplash(slide_data["image_query"])
        if image_url:
            try:
                img_response = requests.get(image_url, timeout=10)
                if img_response.status_code == 200:
                    img_stream = BytesIO(img_response.content)
                    
                    # Add image with frame
                    left = Inches(6.2)
                    top = Inches(1.8)
                    width = Inches(3.3)
                    
                    # Add image
                    pic = slide.shapes.add_picture(img_stream, left, top, width=width)
                    
                    # Add border
                    border = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, left, top, width, pic.height
                    )
                    border.fill.background()
                    border.line.color.rgb = colors["secondary"]
                    border.line.width = Pt(2)
                    
            except Exception as e:
                st.warning(f"Could not add image: {e}")
    
    # Add page number
    if include_page_numbers:
        left = Inches(9)
        top = Inches(6.5)
        width = Inches(0.7)
        height = Inches(0.3)
        
        page_box = slide.shapes.add_textbox(left, top, width, height)
        page_frame = page_box.text_frame
        p = page_frame.paragraphs[0]
        p.text = f"{slide_number}/{total_slides}"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(128, 128, 128)
        p.alignment = PP_ALIGN.RIGHT

def create_conclusion_slide(prs, colors, slide_number, total_slides):
    """Create a beautiful conclusion slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors["primary"]
    
    # Add decorative circles
    for i in range(3):
        left = Inches(i * 3 - 1)
        top = Inches(2)
        width = Inches(4)
        height = Inches(4)
        
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, left, top, width, height
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = colors["secondary"]
        circle.fill.transparency = 0.8
        circle.line.fill.background()
    
    # Add thank you message
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(2)
    
    thank_box = slide.shapes.add_textbox(left, top, width, height)
    thank_frame = thank_box.text_frame
    p = thank_frame.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Add questions text
    left = Inches(1)
    top = Inches(4.5)
    width = Inches(8)
    height = Inches(1)
    
    q_box = slide.shapes.add_textbox(left, top, width, height)
    q_frame = q_box.text_frame
    p = q_frame.paragraphs[0]
    p.text = "Questions?"
    p.font.size = Pt(40)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    p.font.italic = True
    
    # Add page number
    if include_page_numbers:
        left = Inches(9)
        top = Inches(6.5)
        width = Inches(0.7)
        height = Inches(0.3)
        
        page_box = slide.shapes.add_textbox(left, top, width, height)
        page_frame = page_box.text_frame
        p = page_frame.paragraphs[0]
        p.text = f"{slide_number}/{total_slides}"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.RIGHT

def create_presentation(content: dict) -> bytes:
    """Create a beautiful PowerPoint presentation from the generated content."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    total_slides = len(content.get("slides", [])) + 1  # +1 for title slide
    
    # Title slide
    create_beautiful_title_slide(
        prs,
        content.get("title", "Presentation"),
        content.get("subtitle", "An AI Generated Presentation"),
        colors
    )
    
    # Content slides
    for i, slide_data in enumerate(content.get("slides", []), start=2):
        create_beautiful_content_slide(
            prs, slide_data, colors, i, total_slides
        )
    
    # Save to bytes
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()

# Generate presentation
if generate_btn:
    if not topic:
        st.error("Please enter a presentation topic")
    elif not GEMINI_API_KEY:
        st.error("Please configure GEMINI_API_KEY in your .env file")
    elif not available_models:
        st.error("No models available. Please check your API key.")
    else:
        with progress_placeholder.container():
            st.info("🔄 Generating presentation content...")
        
        # Generate content with selected model
        content = generate_presentation_content(topic, num_slides, presentation_style, selected_model)
        
        if content:
            with progress_placeholder.container():
                st.info("🎨 Creating beautiful PowerPoint presentation...")
            
            # Create presentation
            pptx_bytes = create_presentation(content)
            
            with progress_placeholder.container():
                st.success("✅ Beautiful presentation generated successfully!")
            
            # Download button with styling
            col1, col2, col3 = st.columns([1, 2, 1])
            with col2:
                st.download_button(
                    label="📥 Download Beautiful Presentation",
                    data=pptx_bytes,
                    file_name=f"{topic.replace(' ', '_')}_presentation.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True
                )
            
            # Display generated content preview
            with st.expander("📋 View Generated Content"):
                st.json(content)
                
            # Success message
            st.balloons()
        else:
            st.error("Failed to generate presentation. Please try again.")